"""
碳排放核算工具 - 樱桃和苹果
基于GHG Protocol和IPCC 2006标准
"""

import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

# 页面配置
st.set_page_config(
    page_title="碳排放核算工具 - 樱桃和苹果",
    page_icon="🍒",
    layout="wide"
)

# CSS样式
st.markdown("""
<style>
    .main-header {
        font-size: 2.5rem; 
        font-weight: 700; 
        background: linear-gradient(120deg, #e74c3c 0%, #c0392b 50%, #27ae60 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        text-align: center; 
        padding: 1rem 0;
    }
    .sub-header {font-size: 1.2rem; color: #64748b; text-align: center; padding-bottom: 2rem;}
</style>
""", unsafe_allow_html=True)

# 初始化排放因子库
if 'emission_factors' not in st.session_state:
    st.session_state.emission_factors = {
        "固定燃烧-天然气": {"factor": 2.1622, "unit": "kgCO2/m3", "ghg_type": "CO2"},
        "固定燃烧-煤炭": {"factor": 2.38, "unit": "kgCO2/kg", "ghg_type": "CO2"},
        "固定燃烧-柴油": {"factor": 3.0959, "unit": "kgCO2/kg", "ghg_type": "CO2"},
        "固定燃烧-汽油": {"factor": 2.9251, "unit": "kgCO2/kg", "ghg_type": "CO2"},
        "移动燃烧-汽油": {"factor": 2.9251, "unit": "kgCO2/kg", "ghg_type": "CO2"},
        "移动燃烧-柴油": {"factor": 3.0959, "unit": "kgCO2/kg", "ghg_type": "CO2"},
        "工艺排放-丙烷": {"factor": 2.9761, "unit": "kgCO2/kg", "ghg_type": "CO2"},
        "工艺排放-二氧化碳": {"factor": 1.0, "unit": "kgCO2/kg", "ghg_type": "CO2"},
        "无组织排放-R410A": {"factor": 2088, "unit": "kgCO2e/kg", "ghg_type": "HFCs"},
        "无组织排放-R32": {"factor": 675, "unit": "kgCO2e/kg", "ghg_type": "HFCs"},
        "无组织排放-甲烷(化粪池)": {"factor": 22.4, "unit": "kgCO2e/kgBOD", "ghg_type": "CH4"},
        "外购电力-全国平均": {"factor": 0.5703, "unit": "kgCO2/kWh", "ghg_type": "CO2"},
        "外购电力-华北区域": {"factor": 0.8843, "unit": "kgCO2/kWh", "ghg_type": "CO2"},
        "外购电力-华东区域": {"factor": 0.7035, "unit": "kgCO2/kWh", "ghg_type": "CO2"},
        "外购热力-蒸汽": {"factor": 110, "unit": "kgCO2/GJ", "ghg_type": "CO2"},
    }

if 'uploaded_data' not in st.session_state:
    st.session_state.uploaded_data = None
if 'matched_data' not in st.session_state:
    st.session_state.matched_data = None
if 'edited_data' not in st.session_state:
    st.session_state.edited_data = None
if 'calculation_done' not in st.session_state:
    st.session_state.calculation_done = False

# 侧边栏
with st.sidebar:
    st.title("🔧 排放因子管理")
    st.markdown("---")
    
    total_factors = len(st.session_state.emission_factors)
    st.metric("因子总数", total_factors)
    
    with st.expander("📚 查看因子库", expanded=False):
        factor_df = pd.DataFrame([
            {'排放源': k, '排放因子': v['factor'], '单位': v['unit'], '气体': v['ghg_type']}
            for k, v in st.session_state.emission_factors.items()
        ])
        st.dataframe(factor_df, use_container_width=True, height=300)
    
    st.markdown("---")
    st.subheader("➕ 添加排放因子")
    
    new_name = st.text_input("排放源名称", placeholder="例：固定燃烧-生物质", key="new_factor_name")
    col_a, col_b = st.columns(2)
    with col_a:
        new_factor = st.number_input("排放因子", min_value=0.0, step=0.01, format="%.4f", key="new_factor_value")
    with col_b:
        new_unit = st.text_input("单位", placeholder="kgCO2/kg", key="new_factor_unit")
    new_ghg = st.selectbox("温室气体", ["CO2", "CH4", "N2O", "HFCs"], key="new_factor_ghg")
    
    if st.button("✅ 添加到因子库", use_container_width=True):
        if new_name and new_factor > 0 and new_unit:
            st.session_state.emission_factors[new_name] = {
                "factor": new_factor, 
                "unit": new_unit, 
                "ghg_type": new_ghg
            }
            st.success(f"✅ 已添加: {new_name}")
            st.rerun()
        else:
            st.error("⚠️ 请填写完整信息且排放因子必须大于0")

# 主界面
st.markdown('<p class="main-header">🍒🍎 碳排放核算工具 - 樱桃和苹果</p>', unsafe_allow_html=True)
st.markdown('<p class="sub-header">基于 GHG Protocol 和 IPCC 2006 标准 | 支持在线编辑和公式关联</p>', unsafe_allow_html=True)

# 创建模板
def create_template():
    data = {
        '类别': ['范围一：直接温室气体排放']*4 + ['范围二：间接温室气体排放']*2,
        '子类别': ['1.1 固定燃烧', '1.2 移动燃烧', '1.3 工艺排放', '1.4 无组织排放', '2.1 外购电力', '2.2 外购热力'],
        '排放源': ['天然气', '汽油', '丙烷', 'R410A', '外购市政电', '蒸汽'],
        '设施/过程': ['燃气锅炉', '公务车', '焊接', '空调', '用电', '供暖设备'],
        '活动数据': [1239138, 11010, 792, 3.15, 1500000, 500],
        '计量单位': ['m³', 'kg', 'kg', 'kg', 'kWh', 'GJ']
    }
    df = pd.DataFrame(data)
    
    output = BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        df.to_excel(writer, index=False, sheet_name='活动数据')
        ws = writer.sheets['活动数据']
        for i, col in enumerate(['A', 'B', 'C', 'D', 'E', 'F']):
            ws.column_dimensions[col].width = 25
        for cell in ws[1]:
            cell.fill = PatternFill(start_color='E74C3C', end_color='E74C3C', fill_type='solid')
            cell.font = Font(color='FFFFFF', bold=True)
            cell.alignment = Alignment(horizontal='center', vertical='center')
    return output.getvalue()

# 步骤1
st.subheader("📥 步骤1: 下载活动数据模板")
col1, col2 = st.columns([3, 1])
with col1:
    st.info("📌 模板包含6列：类别、子类别、排放源、设施/过程、活动数据、计量单位")
with col2:
    st.download_button("📄 下载模板", create_template(), 
                      "碳排放数据模板.xlsx",
                      "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                      use_container_width=True)

st.markdown("---")

# 步骤2
st.subheader("📤 步骤2: 上传并编辑活动数据")
uploaded_file = st.file_uploader("上传Excel文件", type=['xlsx', 'xls'])

if uploaded_file:
    try:
        df = pd.read_excel(uploaded_file)
        required_cols = ['类别', '子类别', '排放源', '设施/过程', '活动数据', '计量单位']
        
        if not all(col in df.columns for col in required_cols):
            st.error(f"❌ 文件格式不正确！必需列：{', '.join(required_cols)}")
        else:
            st.success("✅ 文件上传成功！")
            st.session_state.uploaded_data = df
            
            # 可编辑的数据表格
            st.markdown("#### ✏️ 在线编辑上传的数据（可修改任何单元格）")
            
            edited_upload_df = st.data_editor(
                df, 
                use_container_width=True, 
                height=400,
                num_rows="dynamic",  # 允许添加/删除行
                key="uploaded_data_editor"
            )
            
            st.session_state.uploaded_data = edited_upload_df
            
            # 下载编辑后的数据
            col_dl1, col_dl2, col_dl3 = st.columns([1, 1, 2])
            with col_dl1:
                def export_edited_data():
                    output = BytesIO()
                    with pd.ExcelWriter(output, engine='openpyxl') as writer:
                        edited_upload_df.to_excel(writer, index=False, sheet_name='活动数据')
                    return output.getvalue()
                
                st.download_button(
                    "💾 下载编辑后的数据", 
                    export_edited_data(),
                    f"编辑后数据_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    use_container_width=True
                )
            
            st.markdown("---")
            st.subheader("🔍 步骤3: 排放因子智能匹配")
            
            if st.button("🚀 开始匹配排放因子", type="primary", use_container_width=True):
                result_df = edited_upload_df.copy()
                result_df['建议排放源类型'] = ""
                result_df['排放因子'] = 0.0
                result_df['因子单位'] = ""
                result_df['温室气体类型'] = "CO2"
                result_df['匹配状态'] = ""
                result_df['数据来源'] = ""
                
                for idx, row in result_df.iterrows():
                    subcat = str(row['子类别'])
                    source = str(row['排放源'])
                    
                    if '1.1' in subcat:
                        key = f"固定燃烧-{source}"
                    elif '1.2' in subcat:
                        key = f"移动燃烧-{source}"
                    elif '1.3' in subcat:
                        key = f"工艺排放-{source}"
                    elif '1.4' in subcat:
                        key = f"无组织排放-{source}"
                    elif '2.1' in subcat:
                        key = "外购电力-全国平均" if "电" in source else f"外购电力-{source}"
                    elif '2.2' in subcat:
                        key = f"外购热力-{source}"
                    else:
                        key = ""
                    
                    if key and key in st.session_state.emission_factors:
                        info = st.session_state.emission_factors[key]
                        result_df.at[idx, '建议排放源类型'] = key
                        result_df.at[idx, '排放因子'] = info['factor']
                        result_df.at[idx, '因子单位'] = info['unit']
                        result_df.at[idx, '温室气体类型'] = info['ghg_type']
                        result_df.at[idx, '匹配状态'] = '✅ 已匹配'
                        result_df.at[idx, '数据来源'] = '因子库'
                    else:
                        result_df.at[idx, '建议排放源类型'] = key if key else "未识别"
                        result_df.at[idx, '排放因子'] = 0.0
                        result_df.at[idx, '因子单位'] = '待补充'
                        result_df.at[idx, '温室气体类型'] = 'CO2'
                        result_df.at[idx, '匹配状态'] = '❌ 未匹配'
                        result_df.at[idx, '数据来源'] = '待补充'
                
                st.session_state.matched_data = result_df
                st.success("✅ 匹配完成！请在下方检查并手动调整")
            
            if st.session_state.matched_data is not None:
                st.markdown("#### 📋 匹配结果（支持手动修改任何值）")
                
                matched_df = st.session_state.matched_data
                
                col1, col2, col3 = st.columns(3)
                total = len(matched_df)
                matched = len(matched_df[matched_df['匹配状态'] == '✅ 已匹配'])
                col1.metric("总活动数", total)
                col2.metric("已匹配", matched, delta=f"{matched/total*100:.0f}%")
                col3.metric("未匹配", total - matched)
                
                st.info("💡 提示：您可以直接修改排放因子、温室气体类型等任何字段，即使因子库中不存在该因子")
                
                # 完全可编辑的表格
                edited_matched_df = st.data_editor(
                    matched_df,
                    use_container_width=True,
                    height=500,
                    column_config={
                        "排放因子": st.column_config.NumberColumn(
                            "排放因子",
                            min_value=0.0,
                            format="%.6f",
                            help="可直接修改，支持手动输入"
                        ),
                        "温室气体类型": st.column_config.SelectboxColumn(
                            "温室气体类型",
                            options=["CO2", "CH4", "N2O", "HFCs", "PFCs", "SF6", "NF3"],
                            help="可选择或修改"
                        ),
                        "活动数据": st.column_config.NumberColumn(
                            "活动数据",
                            format="%.2f"
                        )
                    },
                    key="matched_data_editor"
                )
                
                # 标记手动修改的数据
                for idx in edited_matched_df.index:
                    if edited_matched_df.at[idx, '排放因子'] != matched_df.at[idx, '排放因子']:
                        edited_matched_df.at[idx, '数据来源'] = '手动修改'
                        edited_matched_df.at[idx, '匹配状态'] = '✏️ 手动'
                    elif edited_matched_df.at[idx, '排放因子'] > 0:
                        if edited_matched_df.at[idx, '数据来源'] != '手动修改':
                            edited_matched_df.at[idx, '匹配状态'] = '✅ 已匹配'
                
                st.session_state.edited_data = edited_matched_df
                
                # 显示手动修改项
                manual_items = edited_matched_df[edited_matched_df['数据来源'] == '手动修改']
                if not manual_items.empty:
                    with st.expander(f"✏️ 手动修改项 ({len(manual_items)}个)", expanded=True):
                        st.dataframe(manual_items[['设施/过程', '排放源', '排放因子', '因子单位', '温室气体类型']], 
                                   use_container_width=True)
                
                # 下载匹配结果
                col_dl4, col_dl5, col_dl6 = st.columns([1, 1, 2])
                with col_dl4:
                    def export_matched_data():
                        output = BytesIO()
                        with pd.ExcelWriter(output, engine='openpyxl') as writer:
                            edited_matched_df.to_excel(writer, index=False, sheet_name='匹配结果')
                        return output.getvalue()
                    
                    st.download_button(
                        "💾 下载匹配结果", 
                        export_matched_data(),
                        f"匹配结果_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        use_container_width=True
                    )
                
                if st.button("✅ 确认数据，开始计算", type="primary", use_container_width=True):
                    st.session_state.calculation_done = True
                    st.rerun()
    
    except Exception as e:
        st.error(f"❌ 文件读取失败: {str(e)}")
        st.info("请确保文件格式正确")

# 步骤4: 计算和可视化
if st.session_state.calculation_done and st.session_state.edited_data is not None:
    st.markdown("---")
    st.subheader("📊 步骤4: 排放计算结果与分析")
    
    calc_df = st.session_state.edited_data.copy()
    calc_df['排放量(kgCO2e)'] = calc_df['活动数据'] * calc_df['排放因子']
    calc_df['排放量(tCO2e)'] = calc_df['排放量(kgCO2e)'] / 1000
    calc_df['范围'] = calc_df['类别'].apply(lambda x: '范围一' if '直接' in str(x) else '范围二')
    
    total_emission = calc_df['排放量(tCO2e)'].sum()
    scope_summary = calc_df.groupby('范围')['排放量(tCO2e)'].sum()
    scope1 = scope_summary.get('范围一', 0)
    scope2 = scope_summary.get('范围二', 0)
    
    # 汇总卡片
    st.markdown("### 📈 排放汇总")
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.markdown(f"""
        <div style='background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); 
                    padding: 1.5rem; border-radius: 10px; color: white; box-shadow: 0 4px 6px rgba(0,0,0,0.1);'>
            <h3 style='margin:0; font-size: 1.1rem;'>范围一：直接排放</h3>
            <h2 style='margin:0.5rem 0 0 0; font-size: 2.5rem;'>{scope1:.2f}</h2>
            <p style='margin:0; opacity: 0.9;'>tCO₂e | {scope1/total_emission*100 if total_emission > 0 else 0:.1f}%</p>
        </div>
        """, unsafe_allow_html=True)
    
    with col2:
        st.markdown(f"""
        <div style='background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%); 
                    padding: 1.5rem; border-radius: 10px; color: white; box-shadow: 0 4px 6px rgba(0,0,0,0.1);'>
            <h3 style='margin:0; font-size: 1.1rem;'>范围二：间接排放</h3>
            <h2 style='margin:0.5rem 0 0 0; font-size: 2.5rem;'>{scope2:.2f}</h2>
            <p style='margin:0; opacity: 0.9;'>tCO₂e | {scope2/total_emission*100 if total_emission > 0 else 0:.1f}%</p>
        </div>
        """, unsafe_allow_html=True)
    
    with col3:
        st.markdown(f"""
        <div style='background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%); 
                    padding: 1.5rem; border-radius: 10px; color: white; box-shadow: 0 4px 6px rgba(0,0,0,0.1);'>
            <h3 style='margin:0; font-size: 1.1rem;'>排放总量</h3>
            <h2 style='margin:0.5rem 0 0 0; font-size: 2.5rem;'>{total_emission:.2f}</h2>
            <p style='margin:0; opacity: 0.9;'>tCO₂e | 100%</p>
        </div>
        """, unsafe_allow_html=True)
    
    st.markdown("---")
    
    # 可视化标签页
    tab1, tab2, tab3 = st.tabs(["📊 温室气体分析", "🔥 排放源分析", "📑 数据导出"])
    
    with tab1:
        col1, col2 = st.columns(2)
        
        with col1:
            ghg_summary = calc_df.groupby('温室气体类型')['排放量(tCO2e)'].sum().reset_index()
            ghg_summary = ghg_summary.sort_values('排放量(tCO2e)', ascending=False)
            fig = px.pie(ghg_summary, values='排放量(tCO2e)', names='温室气体类型',
                        title='温室气体排放占比', hole=0.4,
                        color_discrete_sequence=px.colors.qualitative.Set3)
            fig.update_traces(textposition='inside', textinfo='percent+label')
            st.plotly_chart(fig, use_container_width=True)
        
        with col2:
            fig2 = px.bar(ghg_summary, x='温室气体类型', y='排放量(tCO2e)',
                         title='各温室气体排放量', text='排放量(tCO2e)',
                         color='排放量(tCO2e)', color_continuous_scale='Blues')
            fig2.update_traces(texttemplate='%{text:.2f}', textposition='outside')
            st.plotly_chart(fig2, use_container_width=True)
    
    with tab2:
        col1, col2 = st.columns(2)
        
        with col1:
            scope_df = pd.DataFrame({'范围': ['范围一', '范围二'], '排放量': [scope1, scope2]})
            fig3 = px.pie(scope_df, values='排放量', names='范围',
                         title='范围一 vs 范围二', hole=0.4,
                         color_discrete_map={'范围一': '#667eea', '范围二': '#f5576c'})
            fig3.update_traces(textposition='inside', textinfo='percent+label+value')
            st.plotly_chart(fig3, use_container_width=True)
        
        with col2:
            subcat = calc_df.groupby('子类别')['排放量(tCO2e)'].sum().reset_index()
            subcat = subcat.sort_values('排放量(tCO2e)', ascending=False)
            fig4 = px.bar(subcat, x='子类别', y='排放量(tCO2e)',
                         title='各子类别排放量', text='排放量(tCO2e)',
                         color='排放量(tCO2e)', color_continuous_scale='Reds')
            fig4.update_traces(texttemplate='%{text:.2f}', textposition='outside')
            fig4.update_layout(xaxis_tickangle=-45)
            st.plotly_chart(fig4, use_container_width=True)
    
    with tab3:
        st.markdown("### 📥 导出计算结果")
        col1, col2 = st.columns(2)
        
        # Excel导出（带公式关联）
        with col1:
            st.markdown("#### 📊 Excel报告（带公式）")
            
            def export_excel_with_formulas():
                output = BytesIO()
                wb = openpyxl.Workbook()
                
                # 工作表1：详细计算（带公式）
                ws1 = wb.active
                ws1.title = "详细计算"
                
                headers = ['类别', '子类别', '排放源', '设施/过程', '活动数据', '计量单位', 
                          '建议排放源类型', '排放因子', '因子单位', '温室气体类型', '数据来源',
                          '排放量(kgCO2e)', '排放量(tCO2e)']
                
                header_fill = PatternFill(start_color='E74C3C', end_color='E74C3C', fill_type='solid')
                header_font = Font(color='FFFFFF', bold=True, size=11)
                
                for col_idx, header in enumerate(headers, 1):
                    cell = ws1.cell(1, col_idx, header)
                    cell.fill = header_fill
                    cell.font = header_font
                    cell.alignment = Alignment(horizontal='center', vertical='center')
                
                thin_border = Border(
                    left=Side(style='thin'), right=Side(style='thin'),
                    top=Side(style='thin'), bottom=Side(style='thin')
                )
                
                for row_idx, (_, row) in enumerate(calc_df.iterrows(), 2):
                    ws1.cell(row_idx, 1, str(row['类别']))
                    ws1.cell(row_idx, 2, str(row['子类别']))
                    ws1.cell(row_idx, 3, str(row['排放源']))
                    ws1.cell(row_idx, 4, str(row['设施/过程']))
                    ws1.cell(row_idx, 5, float(row['活动数据']))
                    ws1.cell(row_idx, 6, str(row['计量单位']))
                    ws1.cell(row_idx, 7, str(row['建议排放源类型']))
                    ws1.cell(row_idx, 8, float(row['排放因子']))
                    ws1.cell(row_idx, 9, str(row['因子单位']))
                    ws1.cell(row_idx, 10, str(row['温室气体类型']))
                    ws1.cell(row_idx, 11, str(row['数据来源']))
                    
                    # 公式：排放量(kgCO2e) = 活动数据 × 排放因子
                    ws1.cell(row_idx, 12).value = f"=E{row_idx}*H{row_idx}"
                    ws1.cell(row_idx, 12).number_format = '0.0000'
                    
                    # 公式：排放量(tCO2e) = 排放量(kgCO2e) / 1000
                    ws1.cell(row_idx, 13).value = f"=L{row_idx}/1000"
                    ws1.cell(row_idx, 13).number_format = '0.0000'
                    
                    # 添加边框
                    for col in range(1, 14):
                        ws1.cell(row_idx, col).border = thin_border
                
                # 设置列宽
                ws1.column_dimensions['A'].width = 28
                ws1.column_dimensions['B'].width = 18
                ws1.column_dimensions['C'].width = 15
                ws1.column_dimensions['D'].width = 20
                ws1.column_dimensions['E'].width = 15
                ws1.column_dimensions['F'].width = 12
                ws1.column_dimensions['G'].width = 22
                ws1.column_dimensions['H'].width = 15
                ws1.column_dimensions['I'].width = 15
                ws1.column_dimensions['J'].width = 15
                ws1.column_dimensions['K'].width = 12
                ws1.column_dimensions['L'].width = 18
                ws1.column_dimensions['M'].width = 18
                
                # 工作表2：排放汇总（带公式）
                ws2 = wb.create_sheet("排放汇总")
                ws2['A1'] = "排放汇总表"
                ws2['A1'].font = Font(size=16, bold=True, color='E74C3C')
                ws2.merge_cells('A1:D1')
                ws2['A1'].alignment = Alignment(horizontal='center')
                
                ws2['A3'] = "范围"
                ws2['B3'] = "排放量(tCO₂e)"
                ws2['C3'] = "占比(%)"
                ws2['D3'] = "备注"
                for cell in [ws2['A3'], ws2['B3'], ws2['C3'], ws2['D3']]:
                    cell.fill = header_fill
                    cell.font = header_font
                    cell.alignment = Alignment(horizontal='center')
                
                last_row = len(calc_df) + 1
                
                ws2['A4'] = "范围一：直接排放"
                ws2['B4'].value = f'=SUMIF(详细计算!$A$2:$A${last_row},"*直接*",详细计算!$M$2:$M${last_row})'
                ws2['B4'].number_format = '0.00'
                ws2['C4'].value = '=IF(B6>0,B4/B6*100,0)'
                ws2['C4'].number_format = '0.00'
                ws2['D4'] = "固定燃烧+移动燃烧+工艺排放+无组织排放"
                
                ws2['A5'] = "范围二：间接排放"
                ws2['B5'].value = f'=SUMIF(详细计算!$A$2:$A${last_row},"*间接*",详细计算!$M$2:$M${last_row})'
                ws2['B5'].number_format = '0.00'
                ws2['C5'].value = '=IF(B6>0,B5/B6*100,0)'
                ws2['C5'].number_format = '0.00'
                ws2['D5'] = "外购电力+外购热力"
                
                ws2['A6'] = "总排放量"
                ws2['A6'].font = Font(bold=True, size=12)
                ws2['B6'].value = '=B4+B5'
                ws2['B6'].number_format = '0.00'
                ws2['B6'].font = Font(bold=True, size=12)
                ws2['C6'] = '100.00'
                ws2['C6'].font = Font(bold=True)
                ws2['D6'] = "企业温室气体排放总量"
                
                ws2.column_dimensions['A'].width = 25
                ws2.column_dimensions['B'].width = 20
                ws2.column_dimensions['C'].width = 15
                ws2.column_dimensions['D'].width = 35
                
                # 工作表3：温室气体分析（带公式）
                ws3 = wb.create_sheet("温室气体分析")
                ws3['A1'] = "温室气体排放分析"
                ws3['A1'].font = Font(size=16, bold=True, color='E74C3C')
                ws3.merge_cells('A1:C1')
                ws3['A1'].alignment = Alignment(horizontal='center')
                
                ws3['A3'] = "温室气体类型"
                ws3['B3'] = "排放量(tCO₂e)"
                ws3['C3'] = "占比(%)"
                for cell in [ws3['A3'], ws3['B3'], ws3['C3']]:
                    cell.fill = header_fill
                    cell.font = header_font
                    cell.alignment = Alignment(horizontal='center')
                
                ghg_types = calc_df['温室气体类型'].unique()
                for idx, ghg in enumerate(ghg_types, 4):
                    ws3.cell(idx, 1, str(ghg))
                    ws3.cell(idx, 2).value = f'=SUMIF(详细计算!$J$2:$J${last_row},"{ghg}",详细计算!$M$2:$M${last_row})'
                    ws3.cell(idx, 2).number_format = '0.0000'
                    ws3.cell(idx, 3).value = f'=IF(排放汇总!$B$6>0,B{idx}/排放汇总!$B$6*100,0)'
                    ws3.cell(idx, 3).number_format = '0.00'
                
                ws3.column_dimensions['A'].width = 20
                ws3.column_dimensions['B'].width = 20
                ws3.column_dimensions['C'].width = 15
                
                # 工作表4：使用说明
                ws4 = wb.create_sheet("使用说明")
                ws4['A1'] = "📖 Excel报告使用说明"
                ws4['A1'].font = Font(size=14, bold=True, color='E74C3C')
                ws4.merge_cells('A1:B1')
                
                instructions = [
                    ["", ""],
                    ["1. 公式说明", ""],
                    ["", "• 排放量(kgCO₂e) = 活动数据 × 排放因子"],
                    ["", "• 排放量(tCO₂e) = 排放量(kgCO₂e) ÷ 1000"],
                    ["", "• 所有汇总数据使用SUMIF公式自动计算"],
                    ["", ""],
                    ["2. 数据可编辑", ""],
                    ["", "• 可直接修改"详细计算"表中的活动数据或排放因子"],
                    ["", "• 修改后所有排放量会自动重新计算"],
                    ["", "• 汇总表和分析表会自动更新"],
                    ["", ""],
                    ["3. 数据来源标注", ""],
                    ["", "• 因子库：来自内置排放因子数据库"],
                    ["", "• 手动修改：用户手动输入或调整的因子"],
                    ["", ""],
                    ["4. 注意事项", ""],
                    ["", "• 请勿删除表头行"],
                    ["", "• 修改数据时请保持数值格式"],
                    ["", "• 建议保存副本后再进行编辑"],
                ]
                
                for row_idx, (col1, col2) in enumerate(instructions, 2):
                    ws4.cell(row_idx, 1, col1)
                    ws4.cell(row_idx, 2, col2)
                    if "说明" in col1:
                        ws4.cell(row_idx, 1).font = Font(bold=True, size=11)
                
                ws4.column_dimensions['A'].width = 20
                ws4.column_dimensions['B'].width = 60
                
                wb.save(output)
                output.seek(0)
                return output.getvalue()
            
            excel_data = export_excel_with_formulas()
            st.download_button(
                "📥 下载Excel报告（带公式关联）", 
                excel_data,
                f"碳排放核算报告_带公式_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True
            )
            st.success("✅ Excel包含4个工作表，所有数据通过公式关联")
        
        # PPT导出（高级简约风格）
        with col2:
            st.markdown("#### 📽️ PPT演示报告（16:9）")
            
            def create_advanced_ppt():
                prs = Presentation()
                prs.slide_width = Inches(16)
                prs.slide_height = Inches(9)
                
                # 配色方案
                color_primary = RGBColor(231, 76, 60)  # 樱桃红
                color_secondary = RGBColor(39, 174, 96)  # 苹果绿
                color_dark = RGBColor(44, 62, 80)
                color_light = RGBColor(236, 240, 241)
                color_accent = RGBColor(52, 152, 219)
                
                # ========== 第1页：封面 ==========
                slide1 = prs.slides.add_slide(prs.slide_layouts[6])
                slide1.background.fill.solid()
                slide1.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
                
                # 顶部装饰条
                top_bar = slide1.shapes.add_shape(1, Inches(0), Inches(0), Inches(16), Inches(0.3))
                top_bar.fill.solid()
                top_bar.fill.fore_color.rgb = color_primary
                top_bar.line.fill.background()
                
                # 主标题
                title_box = slide1.shapes.add_textbox(Inches(2), Inches(2.5), Inches(12), Inches(1.5))
                tf = title_box.text_frame
                tf.text = "企业碳排放核算报告"
                p = tf.paragraphs[0]
                p.font.size = Pt(56)
                p.font.bold = True
                p.font.color.rgb = color_dark
                p.alignment = PP_ALIGN.CENTER
                
                # 副标题
                subtitle_box = slide1.shapes.add_textbox(Inches(2), Inches(4.2), Inches(12), Inches(0.6))
                stf = subtitle_box.text_frame
                stf.text = "CARBON EMISSION ACCOUNTING REPORT"
                sp = stf.paragraphs[0]
                sp.font.size = Pt(20)
                sp.font.color.rgb = RGBColor(127, 140, 141)
                sp.alignment = PP_ALIGN.CENTER
                
                # 关键数据圆形
                circle = slide1.shapes.add_shape(9, Inches(6.5), Inches(5.5), Inches(3), Inches(3))
                circle.fill.solid()
                circle.fill.fore_color.rgb = color_primary
                circle.line.fill.background()
                
                data_box = slide1.shapes.add_textbox(Inches(6.5), Inches(6.3), Inches(3), Inches(1.5))
                dtf = data_box.text_frame
                dtf.text = f"{total_emission:.1f}\ntCO₂e"
                dtf.paragraphs[0].font.size = Pt(36)
                dtf.paragraphs[0].font.bold = True
                dtf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                dtf.paragraphs[0].alignment = PP_ALIGN.CENTER
                dtf.paragraphs[0].line_spacing = 0.9
                
                # 日期
                date_box = slide1.shapes.add_textbox(Inches(2), Inches(7.8), Inches(12), Inches(0.5))
                date_tf = date_box.text_frame
                date_tf.text = f"{pd.Timestamp.now().strftime('%Y年%m月%d日')} | 基于GHG Protocol & IPCC 2006标准"
                date_p = date_tf.paragraphs[0]
                date_p.font.size = Pt(16)
                date_p.font.color.rgb = RGBColor(149, 165, 166)
                date_p.alignment = PP_ALIGN.CENTER
                
                # ========== 第2页：核心发现 ==========
                slide2 = prs.slides.add_slide(prs.slide_layouts[6])
                slide2.background.fill.solid()
                slide2.background.fill.fore_color.rgb = color_light
                
                # 标题栏
                title_bar = slide2.shapes.add_shape(1, Inches(0), Inches(0), Inches(16), Inches(1.2))
                title_bar.fill.solid()
                title_bar.fill.fore_color.rgb = RGBColor(255, 255, 255)
                title_bar.line.fill.background()
                
                title2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(14.4), Inches(0.6))
                tf2 = title2.text_frame
                tf2.text = "01 | 核心发现与数据概览"
                p2 = tf2.paragraphs[0]
                p2.font.size = Pt(32)
                p2.font.bold = True
                p2.font.color.rgb = color_dark
                
                # 三个数据卡片
                cards_data = [
                    ("范围一：直接排放", scope1, scope1/total_emission*100 if total_emission > 0 else 0, color_accent),
                    ("范围二：间接排放", scope2, scope2/total_emission*100 if total_emission > 0 else 0, color_primary),
                    ("总排放量", total_emission, 100, color_secondary)
                ]
                
                x_positions = [1.5, 6, 10.5]
                for i, (label, value, pct, color) in enumerate(cards_data):
                    # 卡片背景
                    card = slide2.shapes.add_shape(1, Inches(x_positions[i]), Inches(2), Inches(4), Inches(4))
                    card.fill.solid()
                    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    card.line.color.rgb = color
                    card.line.width = Pt(2)
                    
                    # 标签
                    label_box = slide2.shapes.add_textbox(Inches(x_positions[i] + 0.3), Inches(2.4), Inches(3.4), Inches(0.6))
                    ltf = label_box.text_frame
                    ltf.text = label
                    ltf.paragraphs[0].font.size = Pt(16)
                    ltf.paragraphs[0].font.color.rgb = color_dark
                    ltf.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                    # 数值
                    value_box = slide2.shapes.add_textbox(Inches(x_positions[i] + 0.3), Inches(3.2), Inches(3.4), Inches(1))
                    vtf = value_box.text_frame
                    vtf.text = f"{value:.2f}"
                    vtf.paragraphs[0].font.size = Pt(42)
                    vtf.paragraphs[0].font.bold = True
                    vtf.paragraphs[0].font.color.rgb = color
                    vtf.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                    # 单位
                    unit_box = slide2.shapes.add_textbox(Inches(x_positions[i] + 0.3), Inches(4.2), Inches(3.4), Inches(0.4))
                    utf = unit_box.text_frame
                    utf.text = f"tCO₂e"
                    utf.paragraphs[0].font.size = Pt(14)
                    utf.paragraphs[0].font.color.rgb = RGBColor(127, 140, 141)
                    utf.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                    # 占比
                    pct_box = slide2.shapes.add_textbox(Inches(x_positions[i] + 0.3), Inches(4.8), Inches(3.4), Inches(0.6))
                    ptf = pct_box.text_frame
                    ptf.text = f"{pct:.1f}%"
                    ptf.paragraphs[0].font.size = Pt(24)
                    ptf.paragraphs[0].font.bold = True
                    ptf.paragraphs[0].font.color.rgb = color
                    ptf.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # 底部说明
                note_box = slide2.shapes.add_textbox(Inches(1.5), Inches(6.8), Inches(13), Inches(1.2))
                ntf = note_box.text_frame
                main_scope = "范围一" if scope1 > scope2 else "范围二"
                ntf.text = f"💡 关键洞察：企业{main_scope}排放占主导地位（{max(scope1, scope2)/total_emission*100:.1f}%），表明{'直接生产活动' if main_scope == '范围一' else '外购能源消耗'}是主要排放来源。\n建议优先关注{main_scope}的减排机会，可实现最大减排效益。"
                for p in ntf.paragraphs:
                    p.font.size = Pt(14)
                    p.font.color.rgb = color_dark
                    p.line_spacing = 1.4
                
                # ========== 第3页：排放结构分析 ==========
                slide3 = prs.slides.add_slide(prs.slide_layouts[6])
                slide3.background.fill.solid()
                slide3.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
                
                title_bar3 = slide3.shapes.add_shape(1, Inches(0), Inches(0), Inches(16), Inches(1.2))
                title_bar3.fill.solid()
                title_bar3.fill.fore_color.rgb = color_light
                title_bar3.line.fill.background()
                
                title3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(14.4), Inches(0.6))
                tf3 = title3.text_frame
                tf3.text = "02 | 排放结构深度分析"
                p3 = tf3.paragraphs[0]
                p3.font.size = Pt(32)
                p3.font.bold = True
                p3.font.color.rgb = color_dark
                
                # 左侧：排放源表格
                subcat_data = calc_df.groupby('子类别')['排放量(tCO2e)'].sum().reset_index()
                subcat_data = subcat_data.sort_values('排放量(tCO2e)', ascending=False).head(6)
                
                table_title = slide3.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(6), Inches(0.5))
                ttf = table_title.text_frame
                ttf.text = "各子类别排放量明细"
                ttf.paragraphs[0].font.size = Pt(18)
                ttf.paragraphs[0].font.bold = True
                ttf.paragraphs[0].font.color.rgb = color_dark
                
                rows = len(subcat_data) + 1
                table = slide3.shapes.add_table(rows, 3, Inches(1.2), Inches(2.5), Inches(6), Inches(4.5)).table
                
                headers = ['排放源类别', '排放量(tCO₂e)', '占比(%)']
                for i, h in enumerate(headers):
                    cell = table.cell(0, i)
                    cell.text = h
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = color_primary
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.text_frame.paragraphs[0].font.size = Pt(14)
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                for idx, row in subcat_data.iterrows():
                    row_idx = list(subcat_data.index).index(idx) + 1
                    table.cell(row_idx, 0).text = str(row['子类别'])
                    table.cell(row_idx, 1).text = f"{row['排放量(tCO2e)']:.2f}"
                    table.cell(row_idx, 2).text = f"{row['排放量(tCO2e)']/total_emission*100:.1f}%"
                    
                    for col in range(3):
                        cell = table.cell(row_idx, col)
                        cell.text_frame.paragraphs[0].font.size = Pt(12)
                        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        # 高亮最大值
                        if row_idx == 1:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(255, 243, 224)
                
                # 右侧：分析文本
                analysis_box = slide3.shapes.add_textbox(Inches(8), Inches(1.8), Inches(7), Inches(5.7))
                atf = analysis_box.text_frame
                
                top_source = subcat_data.iloc[0]
                analysis_text = f"""📊 排放结构特征

【主要排放源】
• {top_source['子类别']}是最大排放源
• 贡献了{top_source['排放量(tCO2e)']/total_emission*100:.1f}%的总排放量
• 排放量达到{top_source['排放量(tCO2e)']:.2f} tCO₂e

【温室气体构成】"""
                
                ghg_data = calc_df.groupby('温室气体类型')['排放量(tCO2e)'].sum().sort_values(ascending=False)
                for ghg, emission in ghg_data.head(3).items():
                    analysis_text += f"\n• {ghg}: {emission:.2f} tCO₂e ({emission/total_emission*100:.1f}%)"
                
                analysis_text += f"""\n\n【排放集中度】
• TOP3排放源占比：{subcat_data.head(3)['排放量(tCO2e)'].sum()/total_emission*100:.1f}%
• 表明排放高度集中，减排应聚焦重点"""
                
                atf.text = analysis_text
                for p in atf.paragraphs:
                    p.font.size = Pt(13)
                    p.font.color.rgb = color_dark
                    p.line_spacing = 1.5
                
                # ========== 第4页：减排路径规划 ==========
                slide4 = prs.slides.add_slide(prs.slide_layouts[6])
                slide4.background.fill.solid()
                slide4.background.fill.fore_color.rgb = color_light
                
                title_bar4 = slide4.shapes.add_shape(1, Inches(0), Inches(0), Inches(16), Inches(1.2))
                title_bar4.fill.solid()
                title_bar4.fill.fore_color.rgb = RGBColor(255, 255, 255)
                title_bar4.line.fill.background()
                
                title4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(14.4), Inches(0.6))
                tf4 = title4.text_frame
                tf4.text = "03 | 减排路径与行动方案"
                p4 = tf4.paragraphs[0]
                p4.font.size = Pt(32)
                p4.font.bold = True
                p4.font.color.rgb = color_dark
                
                # 三阶段减排路径
                phases = [
                    {
                        "title": "短期行动\n（1年内）",
                        "icon": "🎯",
                        "target": "减排10-15%",
                        "actions": [
                            "能效提升：LED照明、变频空调",
                            "设备优化：定期维护保养",
                            "管理措施：节能制度、培训"
                        ],
                        "x": 1.5
                    },
                    {
                        "title": "中期转型\n（1-3年）",
                        "icon": "🔄",
                        "target": "减排25-35%",
                        "actions": [
                            "能源替代：绿色电力证书",
                            "技术升级：高效设备改造",
                            "体系认证：ISO 50001"
                        ],
                        "x": 6
                    },
                    {
                        "title": "长期目标\n（3-5年）",
                        "icon": "🌟",
                        "target": "碳中和",
                        "actions": [
                            "零碳能源：100%可再生能源",
                            "技术创新：CCUS、氢能",
                            "碳抵消：造林、碳汇项目"
                        ],
                        "x": 10.5
                    }
                ]
                
                for phase in phases:
                    # 卡片
                    card = slide4.shapes.add_shape(1, Inches(phase["x"]), Inches(2), Inches(4), Inches(5.2))
                    card.fill.solid()
                    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    card.line.color.rgb = color_primary
                    card.line.width = Pt(1.5)
                    
                    # 图标和标题
                    icon_box = slide4.shapes.add_textbox(Inches(phase["x"] + 0.3), Inches(2.3), Inches(3.4), Inches(0.8))
                    itf = icon_box.text_frame
                    itf.text = f"{phase['icon']} {phase['title']}"
                    itf.paragraphs[0].font.size = Pt(16)
                    itf.paragraphs[0].font.bold = True
                    itf.paragraphs[0].font.color.rgb = color_primary
                    itf.paragraphs[0].alignment = PP_ALIGN.CENTER
                    itf.paragraphs[0].line_spacing = 1.2
                    
                    # 目标
                    target_box = slide4.shapes.add_textbox(Inches(phase["x"] + 0.3), Inches(3.3), Inches(3.4), Inches(0.5))
                    ttf = target_box.text_frame
                    ttf.text = f"目标：{phase['target']}"
                    ttf.paragraphs[0].font.size = Pt(14)
                    ttf.paragraphs[0].font.color.rgb = color_secondary
                    ttf.paragraphs[0].font.bold = True
                    ttf.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                    # 措施列表
                    actions_box = slide4.shapes.add_textbox(Inches(phase["x"] + 0.5), Inches(4.1), Inches(3), Inches(2.8))
                    atf = actions_box.text_frame
                    for action in phase["actions"]:
                        p = atf.add_paragraph() if atf.text else atf.paragraphs[0]
                        p.text = f"• {action}"
                        p.font.size = Pt(11)
                        p.font.color.rgb = color_dark
                        p.line_spacing = 1.4
                        p.space_before = Pt(6)
                
                # 底部建议
                recommendation_box = slide4.shapes.add_textbox(Inches(1.5), Inches(7.5), Inches(13), Inches(1))
                rtf = recommendation_box.text_frame
                
                if main_scope == "范围二":
                    rec_text = "💡 优先建议：企业范围二排放占主导，建议优先采购绿色电力证书（GEC）或签订可再生能源采购协议（VPPA），可快速实现20-30%的减排目标。"